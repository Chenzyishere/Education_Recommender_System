"""
Auto-generate a Chinese graduation-defense PPT for this project.

Usage:
    .\.venv\Scripts\python.exe utils\generate_defense_ppt.py
"""

from __future__ import annotations

import csv
import json
from pathlib import Path
from typing import Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DATA_DIR = ROOT / "data"
DOC_DIR = ROOT / "documents"
RENDERED_DIR = ROOT / "rendered"

METRICS_CSV = DATA_DIR / "logic_metrics_comparison.csv"
REC_JSON = DATA_DIR / "recommendation_simulation.json"
OUTPUT_PPT = DOC_DIR / "kg_sakt_defense_presentation_zh.pptx"


def set_text_style(run, size=24, bold=False, color=(0, 0, 0), font="Microsoft YaHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = RGBColor(*color)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

    t_frame = slide.shapes.title.text_frame
    t_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    set_text_style(t_frame.paragraphs[0].runs[0], size=40, bold=True, color=(24, 63, 110))

    s_frame = slide.placeholders[1].text_frame
    for p in s_frame.paragraphs:
        if p.runs:
            set_text_style(p.runs[0], size=20, color=(70, 70, 70))


def add_bullet_slide(prs: Presentation, title: str, bullets: List[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    title_tf = slide.shapes.title.text_frame
    if title_tf.paragraphs[0].runs:
        set_text_style(title_tf.paragraphs[0].runs[0], size=32, bold=True, color=(24, 63, 110))

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for idx, line in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        if p.runs:
            set_text_style(p.runs[0], size=20, color=(30, 30, 30))


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    note: str = "",
    left=0.7,
    top=1.7,
    width=12.0,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title
    title_tf = slide.shapes.title.text_frame
    if title_tf.paragraphs[0].runs:
        set_text_style(title_tf.paragraphs[0].runs[0], size=30, bold=True, color=(24, 63, 110))

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))
    else:
        box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(10), Inches(1))
        box.text_frame.text = f"未找到图片: {image_path.name}"
        if box.text_frame.paragraphs[0].runs:
            set_text_style(box.text_frame.paragraphs[0].runs[0], size=20, color=(180, 50, 50))

    if note:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.6))
        box.text_frame.text = note
        p = box.text_frame.paragraphs[0]
        if p.runs:
            set_text_style(p.runs[0], size=14, color=(80, 80, 80))


def read_metrics() -> List[Dict[str, str]]:
    with METRICS_CSV.open("r", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def read_recommendations() -> Dict:
    with REC_JSON.open("r", encoding="utf-8") as f:
        return json.load(f)


def fmt(v: str) -> str:
    if v is None or v == "":
        return "N/A"
    return f"{float(v):.4f}"


def add_metrics_table_slide(prs: Presentation, rows: List[Dict[str, str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "实验结果汇总（核心指标）"
    if slide.shapes.title.text_frame.paragraphs[0].runs:
        set_text_style(slide.shapes.title.text_frame.paragraphs[0].runs[0], size=30, bold=True, color=(24, 63, 110))

    headers = ["模型", "BestEp", "AUC↑", "RMSE↓", "Path↑", "PVR↓", "APC↑", "VS↓", "RDC↑"]
    n_rows = len(rows) + 1
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(0.4),
        Inches(1.4),
        Inches(12.5),
        Inches(4.6),
    )
    table = table_shape.table

    for c, h in enumerate(headers):
        table.cell(0, c).text = h
        p = table.cell(0, c).text_frame.paragraphs[0]
        if p.runs:
            set_text_style(p.runs[0], size=13, bold=True, color=(255, 255, 255))
        table.cell(0, c).fill.solid()
        table.cell(0, c).fill.fore_color.rgb = RGBColor(24, 63, 110)

    for r, row in enumerate(rows, start=1):
        vals = [
            row["Model"],
            row["BestEp"],
            fmt(row["AUC"]),
            fmt(row["RMSE"]),
            fmt(row["Path"]),
            fmt(row["PVR"]),
            fmt(row["APC"]),
            fmt(row["VS"]),
            fmt(row["RDC"]),
        ]
        for c, v in enumerate(vals):
            table.cell(r, c).text = v
            p = table.cell(r, c).text_frame.paragraphs[0]
            if p.runs:
                set_text_style(p.runs[0], size=12, color=(20, 20, 20))


def add_recommendation_case_slide(prs: Presentation, rec: Dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "资源推荐案例分析（导出结果）"
    if slide.shapes.title.text_frame.paragraphs[0].runs:
        set_text_style(slide.shapes.title.text_frame.paragraphs[0].runs[0], size=30, bold=True, color=(24, 63, 110))

    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    # Select three representative personas.
    picked = ["学生A（基础薄弱）", "学生D（高正确率慢节奏）", "学生G（代数偏科）"]
    lines: List[str] = []
    for name in picked:
        if name not in rec:
            continue
        item = rec[name]
        top1 = item["recommendations"][0] if item["recommendations"] else None
        if not top1:
            continue
        lines.append(
            f"{name}：水平{item['level']}，近5次正确率{item['recent_accuracy']:.0%}；"
            f"Top1推荐“{top1['skill_name']}”（分数{top1['score']:.4f}，掌握概率{top1['mastery_prob']:.4f}，认知负荷{top1['cognitive_load']:.4f}）。"
        )

    lines.append("结论：推荐结果体现“先修达标门控 + ZPD区间匹配 + 认知负荷惩罚”的协同作用。")

    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        if p.runs:
            set_text_style(p.runs[0], size=18, color=(30, 30, 30))


def main() -> None:
    DOC_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "基于知识图谱与SAKT的学习资源推荐研究",
        "本科毕业答辩\n计算机科学与技术专业",
    )

    add_bullet_slide(
        prs,
        "研究背景与问题",
        [
            "传统知识追踪模型预测准确，但常忽略教学路径逻辑。",
            "教育推荐需要同时满足：个性化、可达性、可解释性。",
            "核心问题：如何在保持预测性能的同时提升先修逻辑一致性？",
        ],
    )

    add_bullet_slide(
        prs,
        "研究目标与贡献",
        [
            "提出 KG-SAKT 框架：认知预测 + 知识图谱逻辑约束。",
            "构建“预测指标 + 逻辑指标”双维评估体系。",
            "实现推荐阶段策略：掌握学习门控、ZPD匹配、认知负荷惩罚。",
            "输出可解释推荐证据链，支持教学应用。",
        ],
    )

    add_image_slide(
        prs,
        "模型整体逻辑框图",
        RENDERED_DIR / "kg_sakt_logic_diagram.png",
        "图：KG-SAKT 推荐逻辑流程。",
        width=12.0,
    )

    add_bullet_slide(
        prs,
        "数据预处理与图谱构建",
        [
            "数据集：ASSISTments 2009-2010（Skill Builder）。",
            "清洗流程：字段对齐、噪声剔除、序列重构、ID重映射。",
            "图谱构建：三元组 (KP_source, requires, KP_target) + 语义图字典。",
            "计算层：生成 kg_adj_list.json 支撑训练与评估索引。",
        ],
    )

    add_image_slide(
        prs,
        "知识图谱邻接矩阵可视化",
        RENDERED_DIR / "kg_adjacency_matrix.png",
        "图：邻接关系矩阵（用于结构检查与展示）。",
        width=10.8,
        left=1.2,
    )

    add_bullet_slide(
        prs,
        "模型训练与策略",
        [
            "对比模型：Pure-CF、DKT、SAKT、KG-SAKT。",
            "训练损失：L = L_pred + λ_t * L_logic（LogicLambda 退火调度）。",
            "早停策略：基于验证集最优 epoch 选择模型。",
            "推荐阶段：先门控过滤，再融合 ZPD/Readiness/Novelty/Load 排序。",
        ],
    )

    metrics = read_metrics()
    add_metrics_table_slide(prs, metrics)

    add_image_slide(
        prs,
        "逻辑指标对比（柱状图）",
        RENDERED_DIR / "logic_metrics_bar.png",
        "图：Path / PVR / APC / VS / RDC 对比。",
        width=11.4,
        left=0.95,
    )

    add_image_slide(
        prs,
        "逻辑指标对比（雷达图）",
        RENDERED_DIR / "logic_metrics_radar.png",
        "图：多指标综合对比（已做方向归一）。",
        width=11.4,
        left=0.95,
    )

    rec = read_recommendations()
    add_recommendation_case_slide(prs, rec)

    add_bullet_slide(
        prs,
        "结论与展望",
        [
            "KG-SAKT 在保持较强预测能力的同时，显著提升路径逻辑一致性。",
            "推荐结果可输出“前置证据 + 认知负荷”解释链，提升可采纳性。",
            "局限：图谱仍偏轻量规则构建，部分技能语义待完善。",
            "未来：引入专家标注图谱与多模态学习特征，推进在线教学验证。",
        ],
    )

    add_title_slide(prs, "感谢各位老师聆听", "欢迎批评指正")
    prs.save(str(OUTPUT_PPT))
    print(f"[Saved] {OUTPUT_PPT}")


if __name__ == "__main__":
    main()

